# ----------------------------------   1. Imports    -----------------------------
import os, re, math
from hashlib import sha256
from datetime import datetime

import requests
from flask import Flask, Blueprint, request, jsonify, send_from_directory, url_for
from werkzeug.utils import secure_filename
from openai import OpenAI
from dotenv import load_dotenv
from gtts import gTTS
from  Model import  Library

import pdfplumber
import docx
from pptx import Presentation
from openpyxl import load_workbook
from ebooklib import epub, ITEM_DOCUMENT
from sqlalchemy import text

from extensions import mail
from Config import (
    MAIL_SERVER, MAIL_PORT, MAIL_USE_TLS, MAIL_USERNAME, MAIL_PASSWORD,
    MAIL_DEFAULT_SENDER, Session
)
from Config import UPLOADFOLDER
from Model import Books
from Controller import (
    signup_controller, verify_otp_controller, signin_controller, token_required, me_controller,
    forgot_password_controller, reset_password_controller, change_password_controller, _send_reset_email,
    upload_book_controller, get_user_books_controller, update_book_controller,
    delete_book_controller, get_book_controller, get_user_books
)


# ------------------ 2. Base config ------------------
load_dotenv()
app = Flask(__name__)

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
HF_API_KEY = os.getenv("HF_API_KEY")  # Hugging Face key
HF_MODEL = "tiiuae/falcon-7b-instruct"  # free instruct model

app.config['UPLOAD_FOLDER'] = UPLOADFOLDER

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
AUDIO_DIR = os.path.join(BASE_DIR, "static", "audio")
os.makedirs(AUDIO_DIR, exist_ok=True)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

from flask import send_from_directory


@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory("uploads", filename)




ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "pptx", "xlsx", "epub"}

app.config.update(
    MAIL_SERVER=MAIL_SERVER,
    MAIL_PORT=MAIL_PORT,
    MAIL_USE_TLS=MAIL_USE_TLS,
    MAIL_USERNAME=MAIL_USERNAME,
    MAIL_PASSWORD=MAIL_PASSWORD,
    MAIL_DEFAULT_SENDER=MAIL_DEFAULT_SENDER,
)
mail.init_app(app)


# ------------------ 3. Helper functions ------------------
def _slug(s: str) -> str:
    s = (s or "").lower()
    s = re.sub(r"[^a-z0-9]+", "-", s).strip("-")
    return s or "audio"

def _approx_seconds_from_text(txt: str, wpm: int = 150) -> int:
    words = max(1, len((txt or "").split()))
    return math.ceil(words / wpm * 60)

def synthesize_tts_to_file(text: str, title: str, author: str, duration_key: str, lang: str = "en") -> dict:
    digest = sha256((text or "").encode("utf-8")).hexdigest()[:12]
    fname  = f"{_slug(title)}-{_slug(author)}-{duration_key}-{digest}.mp3"
    fpath  = os.path.join(AUDIO_DIR, fname)

    if not os.path.exists(fpath):
        tts = gTTS(text=text, lang=lang)
        tts.save(fpath)

    return {
        "filename": fname,
        "path": fpath,
        "seconds": _approx_seconds_from_text(text)
    }

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(filepath):
    text = ""
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text.strip()

def extract_text_from_docx(filepath):
    doc = docx.Document(filepath)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_txt(filepath):
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_xlsx(filepath):
    wb = load_workbook(filepath, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell])
            if row_text.strip():
                text.append(row_text)
    return "\n".join(text)

def extract_text_from_epub(filepath):
    book = epub.read_epub(filepath)
    text = []
    for item in book.get_items():
        if item.get_type() == ITEM_DOCUMENT:
            try:
                content = item.get_content().decode("utf-8")
                text.append(content)
            except:
                pass
    return "\n".join(text)


# ------------------ 4. Core Routes ------------------

# ---- Auth ----
@app.route('/signup', methods=['POST'])
def signup_route():
    return signup_controller()

@app.route('/verify-otp', methods=['POST'])
def verify_otp_route():
    return verify_otp_controller()

@app.route('/signin', methods=['POST'])
def signin_route():
    return signin_controller()

@app.route("/me", methods=["GET"])
@token_required
def me_route():
    return me_controller()

@app.route('/forgot-password', methods=['POST'])
def forgot_password_route():
    return forgot_password_controller()

@app.route('/reset-password', methods=['POST'])
def reset_password_route():
    return reset_password_controller()

@app.route('/change-password', methods=['POST'])
@token_required
def change_password_route():
    return change_password_controller()


@app.route('/generate-summary', methods=['POST'])
def generate_summary():
    try:
        data = request.get_json() or {}
        title = data.get("title")
        author = data.get("author")
        duration = data.get("duration")

        if not title or not author or not duration:
            return jsonify({"error": "Missing required fields"}), 400

        words_map = {
            "1min": 150,
            "10min": 1500,
            "30min": 4500
        }
        words = words_map.get(duration)
        if not words:
            return jsonify({"error": "Invalid duration"}), 400

        # GPT summary only
        prompt = f"""
        You are a specialized AI assistant acting as a Book Summarizer Bot.
        Your role is that of a "Book Keeper" who has read and learned from a wide range of real books.
        Your task is to produce faithful summaries of books based strictly on their actual content.

        ### Rules & Instructions:
        BOOK NAME = {title}
        AUTHOR = {author}
        WORDS = {words}
        """

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a book summarizer bot."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=words + 200
        )

        summary_text = response.choices[0].message.content.strip()
        if not summary_text:
            return jsonify({"error": "Failed to generate summary"}), 500

        return jsonify({
            "title": title,
            "author": author,
            "duration": duration,
            "target_words": words,
            "summary": summary_text
        }), 200

    except Exception as e:
        print("Error in /generate-summary:", e)
        return jsonify({"error": f"Server error: {str(e)}"}), 500



@app.route('/generate-own-summary', methods=['POST'])
def generate_own_summary():
    try:
        data = request.get_json() or {}
        description = data.get("description")
        duration = data.get("duration")

        if not description or not duration:
            return jsonify({"error": "Missing required fields"}), 400

        words_map = {
            "1min": 150,
            "10min": 1500,
            "30min": 4500
        }
        words = words_map.get(duration)
        if not words:
            return jsonify({"error": "Invalid duration"}), 400

        # GPT prompt for user’s own description
        prompt = f"""
        You are a specialized AI assistant acting as a Book Summarizer Bot.
        You are given a book description written by the user.
        Based strictly on this description, generate a clear and useful summary.

        ### Rules:
        - Do NOT add any external content not hinted at in the description.
        - Expand and compress the summary according to the requested time length.
        - Duration: {duration} (target ~{words} words).
        - The style should be engaging, clear, and faithful to the original description.

        --- Book Description ---
        {description}
        """

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a book summarizer bot."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=words + 200
        )

        summary_text = response.choices[0].message.content.strip()
        if not summary_text:
            return jsonify({"error": "Failed to generate summary"}), 500

        return jsonify({
            "duration": duration,
            "target_words": words,
            "summary": summary_text
        }), 200

    except Exception as e:
        print("Error in /generate-own-summary:", e)
        return jsonify({"error": f"Server error: {str(e)}"}), 500


@app.route('/generate-tts', methods=['POST'])
def generate_tts():
    try:
        data = request.get_json() or {}
        text = data.get("text")
        title = data.get("title", "book")
        author = data.get("author", "unknown")
        duration = data.get("duration", "custom")

        if not text:
            return jsonify({"error": "Missing text for TTS"}), 400

        tts_info = synthesize_tts_to_file(text, title, author, duration)
        audio_url = url_for("serve_audio", filename=tts_info["filename"], _external=True)

        return jsonify({
            "audio_url": audio_url,
            "approx_audio_seconds": tts_info["seconds"],
        }), 200

    except Exception as e:
        print("Error in /generate-tts:", e)
        return jsonify({"error": f"Server error: {str(e)}"}), 500


@app.route("/audio/<path:filename>")
def serve_audio(filename):
    return send_from_directory(AUDIO_DIR, filename, mimetype="audio/mpeg", as_attachment=False)

@app.route("/ask-question", methods=["POST"])
def ask_question():
    data = request.get_json() or {}
    title = data.get("title")
    author = data.get("author")
    summary = data.get("summary") or ""
    question = data.get("question")

    if not title or not author or not question:
        return jsonify({"error": "Missing required fields"}), 400

    prompt = f"""
    You are a knowledgeable assistant answering questions about books.

    Book Title: {title}
    Author: {author}

    Book Summary:
    {summary}

    Question: {question}

    ### Instructions:
    - Provide a clear and accurate answer.
    - Your response must be **4–5 sentences long** (not just one line).
    - Base your response on the given summary, and if unclear, explain thoughtfully.
    - Do not just give names — explain with context.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful book assistant who always provides detailed, multi-sentence answers."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500  # increased to allow longer answers
        )

        answer = response.choices[0].message.content.strip()
        return jsonify({"answer": answer}), 200

    except Exception as e:
        return jsonify({"error": f"Failed to get answer: {str(e)}"}), 500


@app.route("/ask-question", methods=["POST"])
def ask_own_question():
    data = request.get_json() or {}
    description = data.get("description")
    duration = data.get("duration")

    if not description or not duration:
        return jsonify({"error": "Missing required fields"}), 400

    prompt = f"""
    You are a knowledgeable assistant answering questions about books.

    description: {description}
    duration: {duration}

   

    ### Instructions:
    - Provide a clear and accurate answer.
    - Your response must be **4–5 sentences long** (not just one line).
    - Base your response on the given summary, and if unclear, explain thoughtfully.
    - Do not just give names — explain with context.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful book assistant who always provides detailed, multi-sentence answers."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500  # increased to allow longer answers
        )

        answer = response.choices[0].message.content.strip()
        return jsonify({"answer": answer}), 200

    except Exception as e:
        return jsonify({"error": f"Failed to get answer: {str(e)}"}), 500


import json

@app.route("/generate-mcqs", methods=["POST"])
def generate_mcqs():
    data = request.get_json() or {}
    title = data.get("title")
    author = data.get("author")
    summary = data.get("summary") or ""

    if not title or not author:
        return jsonify({"error": "Missing required fields"}), 400

    prompt = f"""
    You are a book quiz generator.

    Book Title: {title}
    Author: {author}

    Book Summary:
    {summary}

    ### Task:
    - Generate 5 multiple-choice questions about this book.
    - Each question must have exactly 4 options.
    - Use this JSON format exactly:

    [
      {{
        "question": "string",
        "options": ["string", "string", "string", "string"],
        "correct": "string"
      }}
    ]

    - The field name for the right answer MUST be "correct".
    - Output JSON array only, no extra text.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You generate quiz questions about books."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1000
        )

        output = response.choices[0].message.content.strip()

        try:
            mcqs = json.loads(output)
        except Exception:
            # Try to repair JSON by extracting valid portion
            start = output.find("[")
            end = output.rfind("]") + 1
            mcqs = json.loads(output[start:end])

        return jsonify({"mcqs": mcqs}), 200

    except Exception as e:
        return jsonify({"error": f"Failed to generate MCQs: {str(e)}"}), 500



# ---- Books ----
@app.route('/books/all')
def all_books():
    page = int(request.args.get("page", 1))
    limit = int(request.args.get("limit", 20))
    offset = (page - 1) * limit
    session = Session()
    books = session.execute(text(f"""
        SELECT * 
        FROM Books 
        ORDER BY book_id 
        OFFSET {offset} ROWS FETCH NEXT {limit} ROWS ONLY
    """)).fetchall()
    return jsonify([dict(row._mapping) for row in books])

@app.route('/books/trending')
def trending_books():
    session = Session()
    books = session.execute(text("""
        SELECT TOP 20 * 
        FROM Books 
        WHERE main_category = 'Literature'
        ORDER BY NEWID()
    """)).fetchall()
    return jsonify([dict(row._mapping) for row in books])



@app.route('/books/featured')
def featured_books():
    session = Session()
    books = session.execute(text("""
        SELECT TOP 20 * 
        FROM Books 
        WHERE main_category = 'History'
        ORDER BY NEWID()
    """)).fetchall()
    return jsonify([dict(row._mapping) for row in books])

@app.route("/books/by-category")
def books_by_category():
    main_category = request.args.get("main_category")
    sub_category = request.args.get("sub_category")

    if not main_category or not sub_category:
        return jsonify({"error": "Missing category parameters"}), 400

    db = Session()
    try:
        stmt = text("""
            SELECT * 
            FROM Books
            WHERE main_category = :main_category
              AND sub_category = :sub_category
        """)
        rows = db.execute(stmt, {
            "main_category": main_category,
            "sub_category": sub_category
        }).fetchall()

        books = [
            {
                "book_id": r.book_id,
                "title": r.title,
                "author": r.author,
                "cover_image_url": r.cover_image_url,
                "main_category": r.main_category,
                "sub_category": r.sub_category
            }
            for r in rows
        ]
        return jsonify(books)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        db.close()

@app.route("/books/categories")
def books_categories():
    db = Session()
    try:
        stmt = text("SELECT DISTINCT main_category, sub_category FROM Books")
        rows = db.execute(stmt).fetchall()

        categories = {}
        for r in rows:
            main_cat = r.main_category
            sub_cat = r.sub_category
            if main_cat not in categories:
                categories[main_cat] = []
            if sub_cat and sub_cat not in categories[main_cat]:
                categories[main_cat].append(sub_cat)
        return jsonify(categories)
    except Exception as e:
        print("Books Categories Error:", e)
        # ✅ return empty dict so frontend won’t get "error"
        return jsonify({}), 200
    finally:
        db.close()


@app.route('/upload-book', methods=['POST'])
@token_required
def upload_book_route():
    return upload_book_controller()

@app.route('/user-books/<int:user_id>', methods=['GET'])
def user_books_route(user_id):
    return get_user_books_controller(user_id)

@app.route('/update-book/<int:book_id>', methods=['PUT'])
def update_book_route(book_id):
    return update_book_controller(book_id)

@app.route('/delete-book/<int:book_id>', methods=['DELETE'])
def delete_book_route(book_id):
    return delete_book_controller(book_id)

@app.route('/get-book/<int:book_id>', methods=['GET'])
def get_book_route(book_id):
    return get_book_controller(book_id)

# ---- File extractors ----
@app.route("/extract-text", methods=["POST"])
def extract_text():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "Empty filename"}), 400

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(filepath)

        ext = filename.rsplit(".", 1)[1].lower()
        text = ""

        if ext == "pdf":
            text = extract_text_from_pdf(filepath)
        elif ext == "docx":
            text = extract_text_from_docx(filepath)
        elif ext == "txt":
            text = extract_text_from_txt(filepath)
        elif ext == "pptx":
            text = extract_text_from_pptx(filepath)
        elif ext == "xlsx":
            text = extract_text_from_xlsx(filepath)
        elif ext == "epub":
            text = extract_text_from_epub(filepath)
        else:
            return jsonify({"error": f"Unsupported file format: {ext}"}), 400

        return jsonify({"filename": filename, "text": text[:5000]})

    return jsonify({"error": "File type not allowed"}), 400


@app.route('/append-pdf-to-book', methods=['POST'])
def append_pdf_to_book():
    try:
        data = request.get_json() or {}
        book_id = data.get("book_id")
        pdf_text = data.get("pdf_text")

        if not book_id or not pdf_text:
            return jsonify({"error": "Missing book_id or pdf_text"}), 400

        db = Session()
        book = db.query(Books).filter(Books.book_id == book_id).first()

        if not book:
            return jsonify({"error": "Book not found"}), 404

        if book.description:
            book.description += f"\n\n--- Added from PDF ({datetime.utcnow().date()}) ---\n{pdf_text}"
        else:
            book.description = pdf_text

        db.commit()
        return jsonify({"message": "PDF content appended to book successfully."}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/get-user-books', methods=['GET'])
@token_required
def get_user_books_route():
    return get_user_books()

 # Library related apis
@app.route("/library/add", methods=["POST"])
def add_to_library():
    data = request.get_json() or {}
    user_id = data.get("user_id")
    book_id = data.get("book_id")

    if not user_id or not book_id:
        return jsonify({"error": "Missing user_id or book_id"}), 400

    db = Session()
    try:
        exists = db.query(Library).filter_by(user_id=user_id, book_id=book_id).first()
        if exists:
            return jsonify({"ok": True, "already": True})

        lib = Library(user_id=user_id, book_id=book_id, created_at=datetime.utcnow())
        db.add(lib)
        db.commit()
        return jsonify({"ok": True})
    except Exception as e:
        db.rollback()
        return jsonify({"error": str(e)}), 500
    finally:
        db.close()


@app.route("/library/remove", methods=["POST"])
def remove_from_library():
    data = request.get_json() or {}
    user_id = data.get("user_id")
    book_id = data.get("book_id")

    if not user_id or not book_id:
        return jsonify({"error": "Missing user_id or book_id"}), 400

    db = Session()
    try:
        db.query(Library).filter_by(user_id=user_id, book_id=book_id).delete()
        db.commit()
        return jsonify({"ok": True})
    except Exception as e:
        db.rollback()
        return jsonify({"error": str(e)}), 500
    finally:
        db.close()


@app.route("/library/check", methods=["GET"])
def check_library():
    user_id = request.args.get("user_id", type=int)
    book_id = request.args.get("book_id", type=int)

    if not user_id or not book_id:
        return jsonify({"error": "Missing user_id or book_id"}), 400

    db = Session()
    try:
        exists = db.query(Library).filter_by(user_id=user_id, book_id=book_id).first()
        return jsonify({"saved": bool(exists)})
    finally:
        db.close()

@app.route("/library/list", methods=["GET"])
def list_library():
    user_id = request.args.get("user_id", type=int)
    if not user_id:
        return jsonify({"error": "Missing user_id"}), 400

    db = Session()
    try:
        q = (
            db.query(Books.book_id, Books.title, Books.author, Books.cover_image_url,
                     Books.main_category, Books.sub_category)
            .join(Library, Library.book_id == Books.book_id)
            .filter(Library.user_id == user_id)
            .order_by(Library.created_at.desc())
        )
        items = [{
            "book_id": r.book_id,
            "title": r.title,
            "author": r.author,
            "cover_image_url": r.cover_image_url,
            "main_category": r.main_category,
            "sub_category": r.sub_category
        } for r in q.all()]
        return jsonify({"items": items})
    finally:
        db.close()


@app.route('/upload-book-cover', methods=['POST'])
def upload_book_cover():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOADFOLDER'], filename)
    file.save(filepath)

    # Save just filename in DB (you can modify your Books table insert here)
    # e.g., book.cover_image_url = filename

    file_url = f"http://192.168.18.98:5000/uploads/{filename}"
    return jsonify({"file_url": file_url}), 200


@app.route('/uploads/<filename>')
def uploaded_files(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)


# ------------------ 5. Main entry ------------------
if __name__ == "__main__":
    app.run(host="0.0.0.0", debug=True)


# from datetime import datetime
#
# from flask import Flask,Blueprint,request,jsonify
#
# from extensions import mail
# from Config import MAIL_SERVER, MAIL_PORT, MAIL_USE_TLS, MAIL_USERNAME, MAIL_PASSWORD, MAIL_DEFAULT_SENDER
# from Controller import signup_controller, verify_otp_controller, signin_controller, token_required, me_controller
# from Controller import forgot_password_controller,reset_password_controller,change_password_controller,_send_reset_email
# from Controller import upload_book_controller,get_user_books_controller,update_book_controller,delete_book_controller,get_book_controller,get_user_books
# from Config import Session
# from Model import Books
# from sqlalchemy import text
# # from book_services import get_public_domain_fulltext,_clean_gutenberg,_normalize_ws,_ia_pick_text_url,fetch_google_books,import_books_with_optional_fulltext
#
# from werkzeug.utils import secure_filename
# import pdfplumber
# import docx
# from pptx import Presentation
# from openpyxl import load_workbook
#
# import os
# from flask import Flask, request, jsonify
# from openai import OpenAI
# from dotenv import load_dotenv
# load_dotenv()
#
# from flask import send_file
# from gtts import gTTS
# import io
# import os, re, io, math
# from hashlib import sha256
# from flask import request, jsonify, send_from_directory, url_for
# from gtts import gTTS
#
# BASE_DIR   = os.path.dirname(os.path.abspath(__file__))
# AUDIO_DIR  = os.path.join(BASE_DIR, "static", "audio")
# os.makedirs(AUDIO_DIR, exist_ok=True)
#
#
# app = Flask(__name__)
#
# # Now getenv will work
# client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
#
# from ebooklib import epub, ITEM_DOCUMENT
#
#
#
# app.config.update(
#     MAIL_SERVER=MAIL_SERVER,
#     MAIL_PORT=MAIL_PORT,
#     MAIL_USE_TLS=MAIL_USE_TLS,
#     MAIL_USERNAME=MAIL_USERNAME,
#     MAIL_PASSWORD=MAIL_PASSWORD,
#     MAIL_DEFAULT_SENDER=MAIL_DEFAULT_SENDER,
# )
# mail.init_app(app)
#
# books_bp = Blueprint("books", __name__)
#
#
# UPLOAD_FOLDER = "uploads"
# os.makedirs(UPLOAD_FOLDER, exist_ok=True)
#
# app = Flask(__name__)
# app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
#
# ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "pptx", "xlsx", "epub"}  # add mobi later if needed
#
#
# @app.route('/signup', methods=['POST'])
# def signup_route():
#     return signup_controller()
#
# @app.route('/verify-otp', methods=['POST'])
# def verify_otp_route():
#     return verify_otp_controller()
#
# @app.route('/signin', methods=['POST'])
# def signin_route():
#     return signin_controller()
#
# @app.route("/me", methods=["GET"])
# @token_required
# def me_route():
#     return me_controller()
#
#
# @app.route('/forgot-password', methods=['POST'])
# def forgot_password_route():
#     return forgot_password_controller()
#
# @app.route('/reset-password', methods=['POST'])
# def reset_password_route():
#     return reset_password_controller()
#
# # for now your token_required is in dev-bypass mode; this will still work
# @app.route('/change-password', methods=['POST'])
# @token_required
# def change_password_route():
#     return change_password_controller()
#
# # ----------------------------------------------------------------
# #                         ------------> Scrapping <---------------
# # ----------------- Main API -----------------
# # All Books (paginated for performance)
# @app.route('/books/all')
# def all_books():
#     page = int(request.args.get("page", 1))
#     limit = int(request.args.get("limit", 20))
#     offset = (page - 1) * limit
#
#     session = Session()
#     books = session.execute(text(f"""
#         SELECT *
#         FROM Books
#         ORDER BY book_id
#         OFFSET {offset} ROWS FETCH NEXT {limit} ROWS ONLY
#     """)).fetchall()
#
#     return jsonify([dict(row._mapping) for row in books])
#
# # Trending (Literature by default)
# @app.route('/books/trending')
# def trending_books():
#     session = Session()
#     books = session.execute(text("""
#         SELECT TOP 20 *
#         FROM Books
#         WHERE main_category = 'Literature'
#         ORDER BY NEWID()
#     """)).fetchall()
#
#     return jsonify([dict(row._mapping) for row in books])
#
# # Featured (History by default)
# @app.route('/books/featured')
# def featured_books():
#     session = Session()
#     books = session.execute(text("""
#         SELECT TOP 20 *
#         FROM Books
#         WHERE main_category = 'History'
#         ORDER BY NEWID()
#     """)).fetchall()
#
#     return jsonify([dict(row._mapping) for row in books])
#
# # Categories + Subcategories
#
#
#
#
# @app.route("/books/by-category")
# def books_by_category():
#     main_category = request.args.get("main_category")
#     sub_category = request.args.get("sub_category")
#
#     if not main_category or not sub_category:
#         return jsonify({"error": "Missing category parameters"}), 400
#
#     db = Session()
#     try:
#         stmt = text("""
#             SELECT *
#             FROM Books
#             WHERE main_category = :main_category
#               AND sub_category = :sub_category
#         """)
#         rows = db.execute(stmt, {
#             "main_category": main_category,
#             "sub_category": sub_category
#         }).fetchall()
#
#         books = [
#             {
#                 "book_id": r.book_id,
#                 "title": r.title,
#                 "author": r.author,
#                 "cover_image_url": r.cover_image_url,
#                 "main_category": r.main_category,
#                 "sub_category": r.sub_category
#             }
#             for r in rows
#         ]
#
#         return jsonify(books)
#     except Exception as e:
#         return jsonify({"error": str(e)}), 500
#     finally:
#         db.close()
#
# @app.route("/books/categories")
# def books_categories():
#     db = Session()
#     try:
#         stmt = text("""
#             SELECT DISTINCT main_category, sub_category
#             FROM Books
#         """)
#         rows = db.execute(stmt).fetchall()
#
#         categories = {}
#         for r in rows:
#             main_cat = r.main_category
#             sub_cat = r.sub_category
#             if main_cat not in categories:
#                 categories[main_cat] = []
#             if sub_cat and sub_cat not in categories[main_cat]:
#                 categories[main_cat].append(sub_cat)
#
#         return jsonify(categories)
#     except Exception as e:
#         return jsonify({"error": str(e)}), 500
#     finally:
#         db.close()
#
#
#
#
# # @app.route('/generate-summary', methods=['POST'])
# # def generate_summary():
# #     try:
# #         data = request.get_json()
# #         title = data.get("title")
# #         author = data.get("author")
# #         duration = data.get("duration")  # "1min" | "10min" | "30min"
# #
# #         if not title or not author or not duration:
# #             return jsonify({"error": "Missing required fields"}), 400
# #
# #         # Map duration to word count
# #         words_map = {
# #             "1min": 150,
# #             "10min": 1500,
# #             "30min": 4500
# #         }
# #         words = words_map.get(duration)
# #         if not words:
# #             return jsonify({"error": "Invalid duration"}), 400
# #
# #         # Build the summarization prompt
# #         prompt = f"""
# #         You are a specialized AI assistant acting as a Book Summarizer Bot.
# #         Your role is that of a "Book Keeper" who has read and learned from a wide range of real books.
# #         Your task is to produce faithful summaries of books based strictly on their actual content.
# #
# #         ### Rules & Instructions:
# #         BOOK NAME = {title}
# #         AUTHOR = {author}
# #         WORDS = {words}
# #         """
# #
# #         # Call GPT
# #         response = client.chat.completions.create(
# #             model="gpt-4o-mini",
# #             messages=[
# #                 {"role": "system", "content": "You are a book summarizer bot."},
# #                 {"role": "user", "content": prompt}
# #             ],
# #             max_tokens=words + 200
# #         )
# #
# #         summary_text = response.choices[0].message.content.strip()
# #
# #         return jsonify({
# #             "title": title,
# #             "author": author,
# #             "duration": duration,
# #             "words": words,
# #             "summary": summary_text
# #         })
# #
# #     except Exception as e:
# #         print("Error in /generate-summary:", e)
# #         return jsonify({"error": "Internal server error"}), 500
#
# @app.route('/generate-summary', methods=['POST'])
# def generate_summary():
#     try:
#         data = request.get_json() or {}
#         title = data.get("title")
#         author = data.get("author")
#         duration = data.get("duration")  # "1min" | "10min" | "30min"
#
#         if not title or not author or not duration:
#             return jsonify({"error": "Missing required fields"}), 400
#
#         # Map duration to target length (≈ words)
#         words_map = {
#             "1min": 150,
#             "10min": 1500,
#             "30min": 4500
#         }
#         words = words_map.get(duration)
#         if not words:
#             return jsonify({"error": "Invalid duration"}), 400
#
#         # --- 1) Build summarization prompt ---
#         prompt = f"""
#         You are a specialized AI assistant acting as a Book Summarizer Bot.
#         Your role is that of a "Book Keeper" who has read and learned from a wide range of real books.
#         Your task is to produce faithful summaries of books based strictly on their actual content.
#
#         ### Rules & Instructions:
#         BOOK NAME = {title}
#         AUTHOR = {author}
#         WORDS = {words}
#         """
#
#         # --- 2) Call GPT to get the summary ---
#         response = client.chat.completions.create(
#             model="gpt-4o-mini",
#             messages=[
#                 {"role": "system", "content": "You are a book summarizer bot."},
#                 {"role": "user", "content": prompt}
#             ],
#             max_tokens=words + 200
#         )
#         summary_text = response.choices[0].message.content.strip()
#         if not summary_text:
#             return jsonify({"error": "Failed to generate summary"}), 500
#
#         # --- 3) Generate TTS with gTTS ---
#         tts_info = synthesize_tts_to_file(summary_text, title, author, duration)
#         audio_url = url_for("serve_audio", filename=tts_info["filename"], _external=True)
#
#         # --- 4) Return response ---
#         return jsonify({
#             "title": title,
#             "author": author,
#             "duration": duration,
#             "target_words": words,
#             "approx_audio_seconds": tts_info["seconds"],
#             "summary": summary_text,
#             "audio_url": audio_url
#         }), 200
#
#     except Exception as e:
#         print("Error in /generate-summary:", e)
#         return jsonify({"error": f"Server error: {str(e)}"}), 500
#
#
#
# @app.route("/audio/<path:filename>")
# def serve_audio(filename):
#     return send_from_directory(AUDIO_DIR, filename, mimetype="audio/mpeg", as_attachment=False)
#
#
#
# # ----------------------------- Upload Book  -----------------------------
#
# @app.route('/upload-book', methods=['POST'])
# @token_required  # optional if you want only logged-in users
# def upload_book_route():
#     return upload_book_controller()
#
# # ------------------------------- User get this own book from this api ---------------------------
#
# @app.route('/user-books/<int:user_id>', methods=['GET'])
# def user_books_route(user_id):
#     return get_user_books_controller(user_id)
#
# @app.route('/update-book/<int:book_id>', methods=['PUT'])
# def update_book_route(book_id):
#     return update_book_controller(book_id)
#
# @app.route('/delete-book/<int:book_id>', methods=['DELETE'])
# def delete_book_route(book_id):
#     return delete_book_controller(book_id)
#
#
# @app.route('/get-book/<int:book_id>', methods=['GET'])
# def get_book_route(book_id):
#     return get_book_controller(book_id)
#
#
#
#
# def allowed_file(filename):
#     return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS
#
# # ----------- Extractors -------------
#
# def extract_text_from_pdf(filepath):
#     text = ""
#     with pdfplumber.open(filepath) as pdf:
#         for page in pdf.pages:
#             text += page.extract_text() or ""
#     return text.strip()
#
# def extract_text_from_docx(filepath):
#     doc = docx.Document(filepath)
#     return "\n".join([p.text for p in doc.paragraphs])
#
# def extract_text_from_txt(filepath):
#     with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
#         return f.read()
#
# def extract_text_from_pptx(filepath):
#     prs = Presentation(filepath)
#     text = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text.append(shape.text)
#     return "\n".join(text)
#
# def extract_text_from_xlsx(filepath):
#     wb = load_workbook(filepath, data_only=True)
#     text = []
#     for sheet in wb.sheetnames:
#         ws = wb[sheet]
#         for row in ws.iter_rows(values_only=True):
#             row_text = " ".join([str(cell) for cell in row if cell])
#             if row_text.strip():
#                 text.append(row_text)
#     return "\n".join(text)
#
# def extract_text_from_epub(filepath):
#     book = epub.read_epub(filepath)
#     text = []
#     for item in book.get_items():
#         if item.get_type() == ITEM_DOCUMENT:
#             try:
#                 content = item.get_content().decode("utf-8")
#                 text.append(content)
#             except:
#                 pass
#     return "\n".join(text)
#
# # MOBI handling would need conversion to EPUB first (using kindlegen or Calibre tools).
#
# # ----------- API Route -------------
#
# @app.route("/extract-text", methods=["POST"])
# def extract_text():
#     if "file" not in request.files:
#         return jsonify({"error": "No file uploaded"}), 400
#
#     file = request.files["file"]
#
#     if file.filename == "":
#         return jsonify({"error": "Empty filename"}), 400
#
#     if file and allowed_file(file.filename):
#         filename = secure_filename(file.filename)
#         filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
#         file.save(filepath)
#
#         ext = filename.rsplit(".", 1)[1].lower()
#         text = ""
#
#         if ext == "pdf":
#             text = extract_text_from_pdf(filepath)
#         elif ext == "docx":
#             text = extract_text_from_docx(filepath)
#         elif ext == "txt":
#             text = extract_text_from_txt(filepath)
#         elif ext == "pptx":
#             text = extract_text_from_pptx(filepath)
#         elif ext == "xlsx":
#             text = extract_text_from_xlsx(filepath)
#         elif ext == "epub":
#             text = extract_text_from_epub(filepath)
#         else:
#             return jsonify({"error": f"Unsupported file format: {ext}"}), 400
#
#         return jsonify({"filename": filename, "text": text[:5000]})  # limit to avoid huge responses
#
#     return jsonify({"error": "File type not allowed"}), 400
#
#
#
# @app.route('/append-pdf-to-book', methods=['POST'])
# def append_pdf_to_book():
#     try:
#         data = request.get_json() or {}
#         book_id = data.get("book_id")
#         pdf_text = data.get("pdf_text")
#
#         if not book_id or not pdf_text:
#             return jsonify({"error": "Missing book_id or pdf_text"}), 400
#
#         db = Session()
#         book = db.query(Books).filter(Books.book_id == book_id).first()
#
#
#         if not book:
#             return jsonify({"error": "Book not found"}), 404
#
#         # Append new PDF text to existing description/content
#         if book.description:
#             book.description += f"\n\n--- Added from PDF ({datetime.utcnow().date()}) ---\n{pdf_text}"
#         else:
#             book.description = pdf_text
#
#         db.commit()
#
#         return jsonify({"message": "PDF content appended to book successfully."}), 200
#     except Exception as e:
#         return jsonify({"error": str(e)}), 500
#
# # ----------------------------- Get User Books -----------------------------
# @app.route('/get-user-books', methods=['GET'])
# @token_required  # optional if you want only logged-in users
# def get_user_books_route():
#     return get_user_books()
#
#
#
#
#
# if __name__ == "__main__":
#     app.run(host="0.0.0.0", debug=True)
#
#
#
# # @app.route("/import_books", methods=["POST"])
# # def import_books():
# #     payload = request.get_json(silent=True) or {}
# #     query = payload.get("query", "classics")
# #     max_results = int(payload.get("max_results", 5))
# #
# #     imported = 0
# #     enriched = 0
# #     session = Session()
# #     try:
# #         # From Gutenberg (Gutendex)
# #         guten_url = f"https://gutendex.com/books?search={query}"
# #         g = requests.get(guten_url, timeout=REQUEST_TIMEOUT).json()
# #         for item in g.get("results", [])[:max_results]:
# #             if item.get("copyright") is False:
# #                 title = item.get("title")
# #                 authors = ", ".join([a.get("name") for a in item.get("authors", []) if a.get("name")])
# #                 fmts = item.get("formats", {})
# #                 txt = fmts.get("text/plain; charset=utf-8") or fmts.get("text/plain")
# #                 description = None
# #                 if txt:
# #                     t = requests.get(txt, timeout=REQUEST_TIMEOUT)
# #                     if t.status_code == 200:
# #                         description = _normalize_ws(_clean_gutenberg(t.text))
# #
# #                 if description:
# #                     book = Books(
# #                         title=title,
# #                         author=authors,
# #                         description=description,
# #                         cover_image_url=fmts.get("image/jpeg", ""),
# #                         category="Public Domain",
# #                         created_at=datetime.utcnow()
# #                     )
# #                     session.add(book)
# #                     imported += 1
# #                     enriched += 1
# #
# #         # From Open Library / Internet Archive
# #         ol_url = f"https://openlibrary.org/search.json?q={query}"
# #         o = requests.get(ol_url, timeout=REQUEST_TIMEOUT).json()
# #         for d in o.get("docs", [])[:max_results]:
# #             if d.get("has_fulltext") and d.get("ebook_access") == "public":
# #                 title = d.get("title")
# #                 author = ", ".join(d.get("author_name", []))
# #                 ia_list = d.get("ia") or []
# #                 description = None
# #                 for ia_id in ia_list:
# #                     txt_url = _ia_pick_text_url(ia_id)
# #                     t = requests.get(txt_url, timeout=REQUEST_TIMEOUT)
# #                     if t.status_code == 200 and len(t.text.strip()) > 1000:
# #                         description = _normalize_ws(t.text)
# #                         break
# #
# #                 if description:
# #                     book = Books(
# #                         title=title,
# #                         author=author,
# #                         description=description,
# #                         cover_image_url=f"https://covers.openlibrary.org/b/id/{d.get('cover_i')}-L.jpg" if d.get("cover_i") else "",
# #                         category="Public Domain",
# #                         created_at=datetime.utcnow()
# #                     )
# #                     session.add(book)
# #                     imported += 1
# #                     enriched += 1
# #
# #         session.commit()
# #     except Exception as e:
# #         session.rollback()
# #         raise e
# #     finally:
# #         session.close()
# #
# #     return jsonify({
# #         "status": "success",
# #         "imported": imported,
# #         "fulltext_enriched": enriched,
# #         "source": ["Project Gutenberg", "Open Library / Internet Archive"]
# #     })
#
#
# #                                 Helper functions
#
# def _slug(s: str) -> str:
#     s = (s or "").lower()
#     s = re.sub(r"[^a-z0-9]+", "-", s).strip("-")
#     return s or "audio"
#
# def _approx_seconds_from_text(txt: str, wpm: int = 150) -> int:
#     words = max(1, len((txt or "").split()))
#     return math.ceil(words / wpm * 60)
#
# def synthesize_tts_to_file(text: str, title: str, author: str, duration_key: str, lang: str = "en") -> dict:
#     """
#     Creates/returns a deterministic mp3 filename for the same summary text,
#     so we don't re-generate if it already exists.
#     """
#     digest = sha256((text or "").encode("utf-8")).hexdigest()[:12]
#     fname  = f"{_slug(title)}-{_slug(author)}-{duration_key}-{digest}.mp3"
#     fpath  = os.path.join(AUDIO_DIR, fname)
#
#     if not os.path.exists(fpath):
#         tts = gTTS(text=text, lang=lang)
#         tts.save(fpath)
#
#     return {
#         "filename": fname,
#         "path": fpath,
#         "seconds": _approx_seconds_from_text(text)
#     }
